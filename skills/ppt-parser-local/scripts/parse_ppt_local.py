#!/usr/bin/env python3
"""本地 PPT 解析脚本。

能力：
- 提取主题色
- 提取每页标题
- 按视觉顺序提取文本（支持组合形状与嵌套组合形状）
- 提取图片（独立图片、填充图片、rels 兜底）
- 基于文件 MD5 的 Markdown 缓存输出

输出契约：
- slides[].images[] 仅包含 image_path 和 size，不包含 image_url
"""

from __future__ import annotations

import argparse
import hashlib
import json
import logging
import os
import zipfile
from io import BytesIO
from typing import Any
import xml.etree.ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


LOGGER = logging.getLogger("ppt_parser_local")


def configure_logging(verbose: bool = False) -> None:
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(
        level=level,
        format="%(asctime)s %(levelname)s %(name)s - %(message)s",
    )


def get_file_md5(file_path: str) -> str:
    hash_md5 = hashlib.md5()
    with open(file_path, "rb") as file_handle:
        for chunk in iter(lambda: file_handle.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()


def extract_theme_colors(pptx_path: str) -> dict[str, str]:
    """从 PPTX 中提取主题色。"""
    slot_mapping = {
        "lt1": "Primary Background (Light 1)",
        "dk1": "Primary Text (Dark 1)",
        "lt2": "Secondary Background (Light 2)",
        "dk2": "Secondary Text (Dark 2)",
        "accent1": "Brand Main Color (Accent 1)",
        "accent2": "Accent 2",
        "accent3": "Accent 3",
        "accent4": "Accent 4",
        "accent5": "Accent 5",
        "accent6": "Accent 6",
        "hlink": "Hyperlink",
        "folHlink": "Followed Hyperlink",
    }

    theme_colors: dict[str, str] = {}
    namespaces = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    try:
        with zipfile.ZipFile(pptx_path, "r") as zip_file:
            theme_files = [
                file_name
                for file_name in zip_file.namelist()
                if "theme" in file_name.lower() and file_name.endswith(".xml")
            ]
            if not theme_files:
                return theme_colors

            root = ET.fromstring(zip_file.read(theme_files[0]))
            clr_scheme = root.find(".//a:clrScheme", namespaces)
            if clr_scheme is None:
                return theme_colors

            for child in clr_scheme:
                slot_name = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                slot_description = slot_mapping.get(slot_name, slot_name)

                color_val = None
                srgb = child.find(".//a:srgbClr", namespaces)
                if srgb is not None:
                    color_val = srgb.get("val")

                if color_val is None:
                    sys_clr = child.find(".//a:sysClr", namespaces)
                    if sys_clr is not None:
                        color_val = sys_clr.get("lastClr")

                if color_val:
                    theme_colors[slot_description] = f"#{color_val}"
    except Exception as exc:  # noqa: BLE001
        LOGGER.error("提取主题色失败: %s", exc)

    return theme_colors


def extract_title_from_slide(slide: Any) -> tuple[str, Any | None]:
    """从幻灯片中提取标题，返回 (title_text, title_shape)。"""
    for shape in slide.shapes:
        try:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    if shape.has_text_frame:
                        title_text = shape.text_frame.text.strip()
                        if title_text:
                            return title_text, shape
        except Exception as exc:  # noqa: BLE001
            LOGGER.debug("占位符标题识别失败: %s", exc)

    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text_frame.text.strip()
            if title_text:
                return title_text, slide.shapes.title
    except Exception as exc:  # noqa: BLE001
        LOGGER.debug("slide.shapes.title 标题识别失败: %s", exc)

    text_shapes: list[tuple[float, str, Any]] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text_frame.text.strip()
        if not text or len(text) >= 100:
            continue

        try:
            top = shape.top
            left = shape.left
            if top < 2_500_000:
                distance = left + top * 0.5
                text_shapes.append((distance, text, shape))
        except Exception as exc:  # noqa: BLE001
            LOGGER.debug("启发式标题识别失败: %s", exc)

    if text_shapes:
        text_shapes.sort(key=lambda item: item[0])
        return text_shapes[0][1], text_shapes[0][2]

    return "", None


def get_shape_position(shape: Any, offset_top: int = 0, offset_left: int = 0) -> tuple[int, int]:
    """获取形状绝对坐标，返回 (top, left)。"""
    try:
        top = shape.top if shape.top is not None else 0
        left = shape.left if shape.left is not None else 0
        return (top + offset_top, left + offset_left)
    except Exception as exc:  # noqa: BLE001
        LOGGER.debug("获取形状坐标失败: %s", exc)
        return (offset_top, offset_left)


def get_group_transform(group_shape: Any) -> dict[str, int]:
    """获取组合形状变换信息。"""
    try:
        grp_sp_pr = group_shape._element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}grpSpPr"
        )
        if grp_sp_pr is None:
            grp_sp_pr = group_shape._element.find(
                ".//{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr"
            )

        if grp_sp_pr is not None:
            xfrm = grp_sp_pr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")
            if xfrm is not None:
                off = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}off")
                ext = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ext")
                ch_off = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chOff")
                ch_ext = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chExt")

                return {
                    "off_x": int(off.get("x", 0)) if off is not None else 0,
                    "off_y": int(off.get("y", 0)) if off is not None else 0,
                    "ext_cx": int(ext.get("cx", 1)) if ext is not None else 1,
                    "ext_cy": int(ext.get("cy", 1)) if ext is not None else 1,
                    "chOff_x": int(ch_off.get("x", 0)) if ch_off is not None else 0,
                    "chOff_y": int(ch_off.get("y", 0)) if ch_off is not None else 0,
                    "chExt_cx": int(ch_ext.get("cx", 1)) if ch_ext is not None else 1,
                    "chExt_cy": int(ch_ext.get("cy", 1)) if ch_ext is not None else 1,
                }
    except Exception as exc:  # noqa: BLE001
        LOGGER.debug("获取组合形状变换信息失败: %s", exc)

    return {
        "off_x": group_shape.left if group_shape.left is not None else 0,
        "off_y": group_shape.top if group_shape.top is not None else 0,
        "ext_cx": 1,
        "ext_cy": 1,
        "chOff_x": 0,
        "chOff_y": 0,
        "chExt_cx": 1,
        "chExt_cy": 1,
    }


def transform_child_position(child_x: int, child_y: int, transform: dict[str, int]) -> tuple[int, int]:
    """将组内子形状坐标转换为幻灯片绝对坐标，返回 (top, left)。"""
    ch_ext_cx = transform["chExt_cx"] if transform["chExt_cx"] != 0 else 1
    ch_ext_cy = transform["chExt_cy"] if transform["chExt_cy"] != 0 else 1

    scale_x = transform["ext_cx"] / ch_ext_cx
    scale_y = transform["ext_cy"] / ch_ext_cy

    absolute_x = transform["off_x"] + (child_x - transform["chOff_x"]) * scale_x
    absolute_y = transform["off_y"] + (child_y - transform["chOff_y"]) * scale_y

    return (int(absolute_y), int(absolute_x))


def extract_text_with_position(
    shape: Any,
    title_text: str | None = None,
    offset_top: int = 0,
    offset_left: int = 0,
) -> list[tuple[str, int, int]]:
    """从普通形状提取文本与坐标。"""
    results: list[tuple[str, int, int]] = []
    top, left = get_shape_position(shape, offset_top, offset_left)

    if shape.has_text_frame:
        paragraphs_text: list[str] = []
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text or (title_text and text == title_text):
                continue
            paragraphs_text.append(text)

        if paragraphs_text:
            results.append(("\n".join(paragraphs_text), top, left))

    if shape.has_table:
        for row_index, row in enumerate(shape.table.rows):
            for col_index, cell in enumerate(row.cells):
                text = cell.text.strip()
                if not text:
                    continue
                cell_top = top + row_index * 100_000
                cell_left = left + col_index * 100_000
                results.append((text, cell_top, cell_left))

    return results


def extract_text_with_transform(
    shape: Any,
    title_text: str | None = None,
    transform: dict[str, int] | None = None,
) -> list[tuple[str, int, int]]:
    """使用组合变换提取子形状文本与绝对坐标。"""
    results: list[tuple[str, int, int]] = []

    try:
        child_x = shape.left if shape.left is not None else 0
        child_y = shape.top if shape.top is not None else 0
        if transform:
            top, left = transform_child_position(child_x, child_y, transform)
        else:
            top, left = child_y, child_x
    except Exception as exc:  # noqa: BLE001
        LOGGER.debug("计算变换坐标失败: %s", exc)
        top, left = 0, 0

    if shape.has_text_frame:
        paragraphs_text: list[str] = []
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text or (title_text and text == title_text):
                continue
            paragraphs_text.append(text)

        if paragraphs_text:
            results.append(("\n".join(paragraphs_text), top, left))

    if shape.has_table:
        for row_index, row in enumerate(shape.table.rows):
            for col_index, cell in enumerate(row.cells):
                text = cell.text.strip()
                if not text:
                    continue
                cell_top = top + row_index * 100_000
                cell_left = left + col_index * 100_000
                results.append((text, cell_top, cell_left))

    return results


def extract_text_from_group(group_shape: Any, title_text: str | None = None) -> list[tuple[str, int, int]]:
    """递归提取组合形状中的文本与绝对坐标。"""
    results: list[tuple[str, int, int]] = []
    transform = get_group_transform(group_shape)

    for sub_shape in group_shape.shapes:
        if sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            results.extend(extract_text_from_group_nested(sub_shape, title_text, transform))
        else:
            results.extend(extract_text_with_transform(sub_shape, title_text, transform))

    return results


def extract_text_from_group_nested(
    group_shape: Any,
    title_text: str | None = None,
    parent_transform: dict[str, int] | None = None,
) -> list[tuple[str, int, int]]:
    """递归处理嵌套组合形状。"""
    results: list[tuple[str, int, int]] = []
    current_transform = get_group_transform(group_shape)

    if parent_transform:
        child_x = current_transform["off_x"]
        child_y = current_transform["off_y"]
        abs_top, abs_left = transform_child_position(child_x, child_y, parent_transform)

        parent_scale_x = parent_transform["ext_cx"] / (
            parent_transform["chExt_cx"] if parent_transform["chExt_cx"] != 0 else 1
        )
        parent_scale_y = parent_transform["ext_cy"] / (
            parent_transform["chExt_cy"] if parent_transform["chExt_cy"] != 0 else 1
        )

        combined_transform = {
            "off_x": abs_left,
            "off_y": abs_top,
            "ext_cx": int(current_transform["ext_cx"] * parent_scale_x),
            "ext_cy": int(current_transform["ext_cy"] * parent_scale_y),
            "chOff_x": current_transform["chOff_x"],
            "chOff_y": current_transform["chOff_y"],
            "chExt_cx": current_transform["chExt_cx"],
            "chExt_cy": current_transform["chExt_cy"],
        }
    else:
        combined_transform = current_transform

    for sub_shape in group_shape.shapes:
        if sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            results.extend(extract_text_from_group_nested(sub_shape, title_text, combined_transform))
        else:
            results.extend(extract_text_with_transform(sub_shape, title_text, combined_transform))

    return results


def normalize_image_ext(content_type: str) -> str:
    ext = content_type.split("/")[-1].lower()
    ext = ext.split("+")[0]
    if ext == "jpeg":
        return "jpg"
    if not ext:
        return "bin"
    return ext


def extract_images_from_slide(slide: Any, slide_index: int, output_dir: str) -> list[dict[str, Any]]:
    """提取幻灯片图片并保存到本地。"""
    image_infos: list[dict[str, Any]] = []
    image_count = 0
    saved_digests: set[str] = set()

    os.makedirs(output_dir, exist_ok=True)

    def save_image(image_blob: bytes, image_ext: str) -> dict[str, Any] | None:
        nonlocal image_count

        digest = hashlib.md5(image_blob).hexdigest()
        if digest in saved_digests:
            return None
        saved_digests.add(digest)

        image_count += 1
        image_filename = f"slide_{slide_index + 1}_image_{image_count}.{image_ext}"
        image_path = os.path.abspath(os.path.join(output_dir, image_filename))

        try:
            with open(image_path, "wb") as file_handle:
                file_handle.write(image_blob)
        except Exception as exc:  # noqa: BLE001
            LOGGER.error("保存图片失败: %s", exc)
            return None

        image_size = [0, 0]
        try:
            with Image.open(BytesIO(image_blob)) as image_file:
                width, height = image_file.size
                image_size = [width, height]
        except Exception as exc:  # noqa: BLE001
            LOGGER.debug("读取图片尺寸失败: %s", exc)

        return {"image_path": image_path, "size": image_size}

    def extract_fill_image(shape: Any) -> dict[str, Any] | None:
        if not hasattr(shape, "fill") or shape.fill.type is None:
            return None

        if shape.fill.type != 6:  # MSO_FILL.PICTURE
            return None

        blip = shape.fill._xPr.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
        )
        if not blip:
            return None

        embed_id = blip.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
        )
        if not embed_id:
            return None

        rel = slide.part.rels.get(embed_id)
        if rel is None:
            return None

        image_part = rel.target_part
        image_blob = image_part.blob
        image_ext = normalize_image_ext(image_part.content_type)
        return save_image(image_blob, image_ext)

    def extract_from_shape(shape: Any) -> None:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                image_info = save_image(image.blob, image.ext)
                if image_info:
                    image_infos.append(image_info)
            except Exception as exc:  # noqa: BLE001
                LOGGER.debug("提取独立图片失败: %s", exc)

        try:
            image_info = extract_fill_image(shape)
            if image_info:
                image_infos.append(image_info)
        except Exception as exc:  # noqa: BLE001
            LOGGER.debug("提取填充图片失败: %s", exc)

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                extract_from_shape(sub_shape)

    for shape in slide.shapes:
        extract_from_shape(shape)

    if image_infos:
        return image_infos

    try:
        for rel in slide.part.rels.values():
            if "image" not in rel.reltype:
                continue

            try:
                image_part = rel.target_part
                image_blob = image_part.blob
                image_ext = normalize_image_ext(image_part.content_type)
                image_info = save_image(image_blob, image_ext)
                if image_info:
                    image_infos.append(image_info)
            except Exception as exc:  # noqa: BLE001
                LOGGER.debug("从 rels 提取图片失败: %s", exc)
    except Exception as exc:  # noqa: BLE001
        LOGGER.debug("遍历 rels 失败: %s", exc)

    return image_infos


def parse_ppt(pptx_path: str, output_dir: str | None = None) -> dict[str, Any]:
    """解析 PPTX 并返回结构化结果。"""
    pptx_path = os.path.abspath(pptx_path)
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"文件不存在: {pptx_path}")

    if output_dir is None:
        output_dir = os.path.join(os.path.dirname(pptx_path), "images")

    os.makedirs(output_dir, exist_ok=True)

    presentation = Presentation(pptx_path)
    theme_colors = extract_theme_colors(pptx_path)

    slides_data: list[dict[str, Any]] = []
    row_threshold = 200_000

    for slide_index, slide in enumerate(presentation.slides):
        slide_title, title_shape = extract_title_from_slide(slide)
        text_with_positions: list[tuple[str, int, int]] = []

        for shape in slide.shapes:
            if title_shape and shape == title_shape:
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                text_with_positions.extend(extract_text_from_group(shape, title_text=slide_title))
            else:
                text_with_positions.extend(extract_text_with_position(shape, title_text=slide_title))

        text_with_positions.sort(key=lambda item: item[1])

        rows: list[list[tuple[str, int, int]]] = []
        current_row: list[tuple[str, int, int]] = []
        current_row_top: int | None = None

        for item in text_with_positions:
            _, top, _ = item
            if current_row_top is None or abs(top - current_row_top) <= row_threshold:
                current_row.append(item)
                if current_row_top is None:
                    current_row_top = top
            else:
                rows.append(current_row)
                current_row = [item]
                current_row_top = top

        if current_row:
            rows.append(current_row)

        sorted_texts: list[str] = []
        for row in rows:
            row.sort(key=lambda item: item[2])
            sorted_texts.extend(item[0] for item in row)

        slide_images = extract_images_from_slide(slide, slide_index, output_dir)

        slide_data = {
            "page": slide_index + 1,
            "title": slide_title,
            "text": sorted_texts,
            "images": slide_images,
        }
        slides_data.append(slide_data)

    return {
        "theme_colors": theme_colors,
        "total_pages": len(presentation.slides),
        "slides": slides_data,
    }


def build_parse_paths(pptx_path: str, output_root: str | None = None) -> dict[str, str]:
    """构建缓存目录与图片目录。"""
    pptx_path = os.path.abspath(pptx_path)
    ppt_name = os.path.splitext(os.path.basename(pptx_path))[0]
    file_md5 = get_file_md5(pptx_path)

    base_dir = os.path.abspath(output_root) if output_root else os.path.dirname(pptx_path)
    parse_dir = os.path.join(base_dir, "pptParseResult", f"{ppt_name}_{file_md5}")
    md_path = os.path.join(parse_dir, "content.md")
    images_dir = os.path.join(parse_dir, "images")

    return {
        "base_dir": base_dir,
        "parse_dir": parse_dir,
        "md_path": md_path,
        "images_dir": images_dir,
    }


def parse_ppt_to_markdown(pptx_path: str, output_root: str | None = None) -> str:
    """解析 PPT 并返回 Markdown；命中缓存时直接读取缓存。"""
    paths = build_parse_paths(pptx_path, output_root)
    md_path = paths["md_path"]

    if os.path.exists(md_path):
        with open(md_path, "r", encoding="utf-8") as file_handle:
            return file_handle.read()

    result = parse_ppt(pptx_path, paths["images_dir"])
    json_content = json.dumps(result, ensure_ascii=False, indent=2)

    md_content = f"""# PPT解析结果

```json
{json_content}
```
"""

    os.makedirs(os.path.dirname(md_path), exist_ok=True)
    with open(md_path, "w", encoding="utf-8") as file_handle:
        file_handle.write(md_content)

    return md_content


def parse_ppt_to_json(pptx_path: str, output_root: str | None = None) -> dict[str, Any]:
    """解析 PPT 并返回结构化 JSON 数据。"""
    paths = build_parse_paths(pptx_path, output_root)
    return parse_ppt(pptx_path, paths["images_dir"])


def build_cli_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="本地 PPT 解析工具")
    parser.add_argument("pptx_path", help="待解析的 .pptx 文件路径")
    parser.add_argument(
        "--mode",
        choices=["json", "markdown"],
        default="markdown",
        help="输出模式：json 或 markdown（默认 markdown）",
    )
    parser.add_argument(
        "--output-dir",
        default=None,
        help="解析结果根目录（默认使用 PPT 所在目录）",
    )
    parser.add_argument(
        "--pretty",
        action="store_true",
        help="json 模式下格式化输出（缩进 2 空格）",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="输出调试日志",
    )
    return parser


def main() -> int:
    parser = build_cli_parser()
    args = parser.parse_args()

    configure_logging(args.verbose)

    try:
        if args.mode == "markdown":
            content = parse_ppt_to_markdown(args.pptx_path, args.output_dir)
            print(content)
            return 0

        result = parse_ppt_to_json(args.pptx_path, args.output_dir)
        if args.pretty:
            print(json.dumps(result, ensure_ascii=False, indent=2))
        else:
            print(json.dumps(result, ensure_ascii=False))
        return 0
    except Exception as exc:  # noqa: BLE001
        LOGGER.error("解析失败: %s", exc)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
